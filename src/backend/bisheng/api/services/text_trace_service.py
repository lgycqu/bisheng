"""
Text Trace Service - 文本溯源服务

提供文本溯源功能，支持精确匹配（Elasticsearch）和语义匹配（Milvus）
"""
import os
import secrets
from enum import Enum
from typing import List, Optional

from loguru import logger
from pydantic import BaseModel, Field

from bisheng.common.dependencies.user_deps import UserPayload
from bisheng.core.cache.redis_manager import get_redis_client, get_redis_client_sync
from bisheng.core.storage.minio.minio_manager import get_minio_storage, get_minio_storage_sync
from bisheng.database.models.role_access import AccessType, RoleAccessDao
from bisheng.knowledge.domain.knowledge_rag import KnowledgeRag
from bisheng.knowledge.domain.models.knowledge import Knowledge, KnowledgeDao
from bisheng.knowledge.domain.models.knowledge_file import KnowledgeFile, KnowledgeFileDao
from bisheng.user.domain.models.user_role import UserRoleDao


class MatchMode(str, Enum):
    """匹配模式枚举"""
    EXACT = "exact"  # 精确匹配（Elasticsearch）
    SEMANTIC = "semantic"  # 语义匹配（Milvus）
    HYBRID = "hybrid"  # 混合匹配


class TextTraceRequest(BaseModel):
    """文本溯源请求"""
    text: str = Field(..., description="待溯源的文本")
    match_mode: MatchMode = Field(default=MatchMode.HYBRID, description="匹配模式")
    top_k: int = Field(default=10, ge=1, le=100, description="返回结果数量")
    threshold: float = Field(default=0.7, ge=0.0, le=1.0, description="匹配阈值")


class MatchResult(BaseModel):
    """匹配结果"""
    document_id: str = Field(..., description="文档ID")
    document_name: str = Field(..., description="文档名称")
    knowledge_base: str = Field(..., description="知识库名称")
    score: float = Field(..., description="匹配分数")
    preview_url: str = Field(default="", description="预览URL")
    matched_text: str = Field(..., description="匹配的文本内容")


class TextTraceResponse(BaseModel):
    """文本溯源响应"""
    matches: List[MatchResult] = Field(default_factory=list, description="匹配结果列表")
    total: int = Field(default=0, description="匹配总数")


class PreviewTokenData(BaseModel):
    """预览Token数据"""
    file_id: int
    user_id: int
    highlight_text: str = Field(default="", max_length=2000)  # 限制高亮文本长度


# Redis key 前缀
PREVIEW_TOKEN_PREFIX = "preview_token:"
# Token 有效期（秒）- 30分钟
PREVIEW_TOKEN_EXPIRATION = 30 * 60


class TextTraceService:
    """文本溯源服务"""

    @classmethod
    async def get_user_knowledge_ids(cls, user_id: int, user_role: List[int]) -> List[int]:
        """
        获取用户有权限的知识库 ID 列表

        Args:
            user_id: 用户ID
            user_role: 用户角色ID列表

        Returns:
            用户有权限访问的知识库ID列表
        """
        # 检查是否是管理员（role_id == 1）
        is_admin = 1 in user_role if user_role else False

        if is_admin:
            # 管理员可以访问所有知识库
            all_knowledge = await KnowledgeDao.aget_all_knowledge()
            return [k.id for k in all_knowledge]

        # 获取用户自己创建的知识库
        user_knowledge = await KnowledgeDao.aget_user_knowledge(user_id)
        knowledge_ids = {k.id for k in user_knowledge}

        # 获取通过角色授权的知识库
        if user_role:
            role_access = await RoleAccessDao.aget_role_access(user_role, AccessType.KNOWLEDGE)
            for access in role_access:
                try:
                    knowledge_ids.add(int(access.third_id))
                except (ValueError, TypeError):
                    continue

        return list(knowledge_ids)

    @classmethod
    async def exact_search(
        cls,
        text: str,
        knowledge_ids: List[int],
        top_k: int,
        threshold: float
    ) -> List[MatchResult]:
        """
        精确匹配搜索（使用 Elasticsearch）

        Args:
            text: 搜索文本
            knowledge_ids: 知识库ID列表
            top_k: 返回结果数量
            threshold: 匹配阈值

        Returns:
            匹配结果列表
        """
        if not knowledge_ids:
            return []

        results = []
        knowledge_list = await KnowledgeDao.aget_list_by_ids(knowledge_ids)
        knowledge_map = {k.id: k for k in knowledge_list}

        for knowledge_id in knowledge_ids:
            knowledge = knowledge_map.get(knowledge_id)
            if not knowledge or not knowledge.index_name:
                continue

            try:
                es_client = KnowledgeRag.init_es_vectorstore_sync(knowledge.index_name)

                # 构建 Elasticsearch 查询
                search_body = {
                    "size": top_k,
                    "query": {
                        "match": {
                            "text": {
                                "query": text,
                                "minimum_should_match": "70%"
                            }
                        }
                    },
                    "_source": ["text", "metadata"]
                }

                response = es_client.client.search(
                    index=knowledge.index_name,
                    body=search_body
                )

                max_score = response["hits"]["max_score"] or 1.0

                for hit in response["hits"]["hits"]:
                    # 归一化分数到 0-1 范围
                    normalized_score = hit["_score"] / max_score if max_score > 0 else 0

                    if normalized_score < threshold:
                        continue

                    metadata = hit["_source"].get("metadata", {})
                    document_id = str(metadata.get("document_id", ""))
                    document_name = metadata.get("document_name", "")
                    matched_text = hit["_source"].get("text", "")

                    # 从拼接后的chunk中分离出原始chunk
                    if matched_text.startswith("{<file_title>"):
                        matched_text = matched_text.split("<paragraph_content>")[-1]
                        matched_text = matched_text.split("</paragraph_content>")[0]

                    results.append(MatchResult(
                        document_id=document_id,
                        document_name=document_name,
                        knowledge_base=knowledge.name,
                        score=round(normalized_score, 4),
                        preview_url="",
                        matched_text=matched_text[:500]  # 限制文本长度
                    ))

            except Exception as e:
                logger.exception(f"Elasticsearch search failed for knowledge {knowledge_id}")
                continue

        # 按分数排序并限制数量
        results.sort(key=lambda x: x.score, reverse=True)
        return results[:top_k]

    @classmethod
    async def semantic_search(
        cls,
        text: str,
        knowledge_ids: List[int],
        top_k: int,
        threshold: float,
        user_id: int
    ) -> List[MatchResult]:
        """
        语义匹配搜索（使用 Milvus）

        Args:
            text: 搜索文本
            knowledge_ids: 知识库ID列表
            top_k: 返回结果数量
            threshold: 匹配阈值
            user_id: 用户ID（用于获取embedding）

        Returns:
            匹配结果列表
        """
        if not knowledge_ids:
            return []

        results = []
        knowledge_list = await KnowledgeDao.aget_list_by_ids(knowledge_ids)
        knowledge_map = {k.id: k for k in knowledge_list}

        for knowledge_id in knowledge_ids:
            knowledge = knowledge_map.get(knowledge_id)
            if not knowledge or not knowledge.collection_name:
                continue

            try:
                milvus_client = await KnowledgeRag.init_knowledge_milvus_vectorstore(
                    invoke_user_id=user_id,
                    knowledge=knowledge
                )

                # 使用 Milvus 进行相似度搜索
                search_results = milvus_client.similarity_search_with_score(
                    query=text,
                    k=top_k
                )

                for doc, score in search_results:
                    # Milvus L2 距离转换为相似度分数（距离越小越相似）
                    # 使用 1 / (1 + distance) 转换，确保分数在 0-1 范围内
                    similarity_score = min(1.0, 1 / (1 + max(0, score)))

                    if similarity_score < threshold:
                        continue

                    metadata = doc.metadata or {}
                    document_id = str(metadata.get("document_id", ""))
                    document_name = metadata.get("document_name", "")
                    matched_text = doc.page_content or ""

                    # 从拼接后的chunk中分离出原始chunk
                    if matched_text.startswith("{<file_title>"):
                        matched_text = matched_text.split("<paragraph_content>")[-1]
                        matched_text = matched_text.split("</paragraph_content>")[0]

                    results.append(MatchResult(
                        document_id=document_id,
                        document_name=document_name,
                        knowledge_base=knowledge.name,
                        score=round(similarity_score, 4),
                        preview_url="",
                        matched_text=matched_text[:500]  # 限制文本长度
                    ))

            except Exception as e:
                logger.exception(f"Milvus search failed for knowledge {knowledge_id}")
                continue

        # 按分数排序并限制数量
        results.sort(key=lambda x: x.score, reverse=True)
        return results[:top_k]

    @classmethod
    async def search(
        cls,
        request: TextTraceRequest,
        user_id: int,
        user_role: List[int]
    ) -> TextTraceResponse:
        """
        主搜索方法，根据 match_mode 调用不同的搜索

        Args:
            request: 搜索请求
            user_id: 用户ID
            user_role: 用户角色列表

        Returns:
            搜索响应
        """
        # 获取用户有权限的知识库
        knowledge_ids = await cls.get_user_knowledge_ids(user_id, user_role)

        if not knowledge_ids:
            return TextTraceResponse(matches=[], total=0)

        results = []

        if request.match_mode == MatchMode.EXACT:
            # 只执行精确匹配
            results = await cls.exact_search(
                text=request.text,
                knowledge_ids=knowledge_ids,
                top_k=request.top_k,
                threshold=request.threshold
            )

        elif request.match_mode == MatchMode.SEMANTIC:
            # 只执行语义匹配
            results = await cls.semantic_search(
                text=request.text,
                knowledge_ids=knowledge_ids,
                top_k=request.top_k,
                threshold=request.threshold,
                user_id=user_id
            )

        elif request.match_mode == MatchMode.HYBRID:
            # 混合模式：先精确匹配，不足 top_k 时补充语义匹配
            exact_results = await cls.exact_search(
                text=request.text,
                knowledge_ids=knowledge_ids,
                top_k=request.top_k,
                threshold=request.threshold
            )

            results = exact_results

            # 如果精确匹配结果不足，补充语义匹配
            if len(results) < request.top_k:
                remaining = request.top_k - len(results)
                semantic_results = await cls.semantic_search(
                    text=request.text,
                    knowledge_ids=knowledge_ids,
                    top_k=remaining + len(results),  # 多取一些用于去重
                    threshold=request.threshold,
                    user_id=user_id
                )

                # 去重：基于 document_id 和 matched_text
                existing_keys = {
                    (r.document_id, r.matched_text[:100])
                    for r in results
                }

                for sr in semantic_results:
                    key = (sr.document_id, sr.matched_text[:100])
                    if key not in existing_keys and len(results) < request.top_k:
                        results.append(sr)
                        existing_keys.add(key)

        # 按分数排序
        results.sort(key=lambda x: x.score, reverse=True)
        results = results[:request.top_k]

        return TextTraceResponse(
            matches=results,
            total=len(results)
        )

    @classmethod
    async def create_preview_token(
        cls,
        file_id: int,
        user_id: int,
        highlight_text: str
    ) -> str:
        """
        创建预览临时 Token（存入 Redis）

        Args:
            file_id: 文件ID
            user_id: 用户ID
            highlight_text: 高亮文本

        Returns:
            生成的 Token
        """
        # 生成安全的随机 Token
        token = secrets.token_urlsafe(32)

        # 构建 Token 数据
        token_data = PreviewTokenData(
            file_id=file_id,
            user_id=user_id,
            highlight_text=highlight_text
        )

        # 存入 Redis
        redis_client = await get_redis_client()
        key = f"{PREVIEW_TOKEN_PREFIX}{token}"
        await redis_client.aset(key, token_data.model_dump(), expiration=PREVIEW_TOKEN_EXPIRATION)

        logger.info(f"Created preview token: {token} for file_id: {file_id}, user_id: {user_id}")

        return token

    @classmethod
    async def validate_preview_token(cls, token: str) -> Optional[PreviewTokenData]:
        """
        验证并消费预览 Token（一次性使用，验证后删除）

        Args:
            token: 预览 Token

        Returns:
            Token 数据，如果无效则返回 None
        """
        redis_client = await get_redis_client()
        key = f"{PREVIEW_TOKEN_PREFIX}{token}"

        # 获取 Token 数据
        token_data = await redis_client.aget(key)

        if not token_data:
            logger.warning(f"Preview token not found or expired: {token}")
            return None

        # 删除 Token（一次性使用）
        await redis_client.adelete(key)

        logger.info(f"Validated and consumed preview token: {token}")

        # 返回 Token 数据
        return PreviewTokenData(**token_data)

    @classmethod
    def create_preview_token_sync(
        cls,
        file_id: int,
        user_id: int,
        highlight_text: str
    ) -> str:
        """
        同步版本：创建预览临时 Token
        """
        # 生成安全的随机 Token
        token = secrets.token_urlsafe(32)

        token_data = PreviewTokenData(
            file_id=file_id,
            user_id=user_id,
            highlight_text=highlight_text
        )

        redis_client = get_redis_client_sync()
        key = f"{PREVIEW_TOKEN_PREFIX}{token}"
        redis_client.set(key, token_data.model_dump(), expiration=PREVIEW_TOKEN_EXPIRATION)

        logger.info(f"Created preview token (sync): {token} for file_id: {file_id}")

        return token

    @classmethod
    def validate_preview_token_sync(cls, token: str) -> Optional[PreviewTokenData]:
        """
        同步版本：验证并消费预览 Token
        """
        redis_client = get_redis_client_sync()
        key = f"{PREVIEW_TOKEN_PREFIX}{token}"

        token_data = redis_client.get(key)

        if not token_data:
            logger.warning(f"Preview token not found or expired (sync): {token}")
            return None

        redis_client.delete(key)

        logger.info(f"Validated and consumed preview token (sync): {token}")

        return PreviewTokenData(**token_data)

    @classmethod
    async def get_file_preview_url(cls, file_id: int) -> str:
        """
        获取文件预览 URL

        Args:
            file_id: 文件ID

        Returns:
            预览 URL
        """
        file_info = await KnowledgeFileDao.query_by_id(file_id)
        if not file_info:
            return ""

        minio_client = await get_minio_storage()

        # 尝试获取预览文件 URL
        if file_info.object_name:
            try:
                url = await minio_client.get_share_link(file_info.object_name)
                return url
            except Exception as e:
                logger.warning(f"Failed to get preview URL for file {file_id}: {e}")

        return ""

    @classmethod
    async def get_file_content_for_preview(cls, file: "KnowledgeFile") -> str:
        """获取文件内容用于预览

        从 MinIO 下载文件并根据文件类型提取文本内容。

        支持的文件格式:
        - PDF: 使用 fitz (PyMuPDF) 提取文本
        - Word (doc, docx): 使用 python-docx 提取文本
        - Excel (xls, xlsx): 使用 pandas 读取并转为字符串
        - PPT (ppt, pptx): 使用 python-pptx 提取文本
        - 文本文件 (txt, md, csv): 直接读取

        Args:
            file: KnowledgeFile 对象

        Returns:
            文件内容字符串，如果解析失败则返回错误提示信息
        """
        if not file or not file.object_name:
            return "文件不存在或文件路径为空"

        # 获取文件扩展名
        file_name = file.file_name or file.object_name
        file_ext = os.path.splitext(file_name)[1].lower()

        try:
            # 从 MinIO 下载文件内容
            minio_client = await get_minio_storage()
            file_bytes = await minio_client.get_object(
                bucket_name=minio_client.bucket,
                object_name=file.object_name
            )

            if not file_bytes:
                return "无法从存储中获取文件内容"

            # 根据文件类型选择解析方式
            if file_ext == '.pdf':
                return cls._extract_pdf_content(file_bytes)
            elif file_ext in ['.doc', '.docx']:
                return cls._extract_docx_content(file_bytes, file_ext)
            elif file_ext in ['.xls', '.xlsx']:
                return cls._extract_excel_content(file_bytes, file_ext)
            elif file_ext in ['.ppt', '.pptx']:
                return cls._extract_pptx_content(file_bytes, file_ext)
            elif file_ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']:
                return cls._extract_text_content(file_bytes)
            else:
                return f"不支持的文件格式: {file_ext}"

        except Exception as e:
            logger.exception(f"Failed to extract content from file {file.id}: {e}")
            return f"文件内容提取失败: {str(e)}"

    @classmethod
    def _extract_pdf_content(cls, file_bytes: bytes) -> str:
        """从 PDF 文件提取文本内容"""
        try:
            import fitz  # PyMuPDF

            # 从字节流打开 PDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text_content = []

            try:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text("text")
                    if text.strip():
                        text_content.append(text)
            finally:
                doc.close()

            return "\n\n".join(text_content) if text_content else "PDF 文件内容为空"

        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return f"PDF 解析失败: {str(e)}"

    @classmethod
    def _extract_docx_content(cls, file_bytes: bytes, file_ext: str) -> str:
        """从 Word 文档提取文本内容"""
        try:
            from io import BytesIO
            from docx import Document

            # .doc 格式需要先转换，这里只支持 .docx
            if file_ext == '.doc':
                return "暂不支持 .doc 格式的内容预览，请转换为 .docx 格式"

            # 从字节流打开文档
            doc = Document(BytesIO(file_bytes))
            text_content = []

            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))

            return "\n\n".join(text_content) if text_content else "Word 文档内容为空"

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return f"Word 文档解析失败: {str(e)}"

    @classmethod
    def _extract_excel_content(cls, file_bytes: bytes, file_ext: str) -> str:
        """从 Excel 文件提取文本内容"""
        try:
            from io import BytesIO
            import pandas as pd

            # 读取 Excel 文件
            excel_file = BytesIO(file_bytes)

            # 根据扩展名选择引擎
            engine = 'openpyxl' if file_ext == '.xlsx' else 'xlrd'

            try:
                # 读取所有工作表
                excel_data = pd.read_excel(excel_file, sheet_name=None, engine=engine, dtype=str)
            except Exception:
                # 如果指定引擎失败，尝试自动检测
                excel_file.seek(0)
                excel_data = pd.read_excel(excel_file, sheet_name=None, dtype=str)

            text_content = []

            for sheet_name, df in excel_data.items():
                if df.empty:
                    continue

                # 添加工作表名称
                text_content.append(f"=== 工作表: {sheet_name} ===")

                # 将 DataFrame 转换为 Markdown 表格格式
                df = df.fillna("")
                table_str = df.to_markdown(index=False)
                if table_str:
                    text_content.append(table_str)

            return "\n\n".join(text_content) if text_content else "Excel 文件内容为空"

        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            return f"Excel 文件解析失败: {str(e)}"

    @classmethod
    def _extract_pptx_content(cls, file_bytes: bytes, file_ext: str) -> str:
        """从 PPT 文件提取文本内容"""
        try:
            from io import BytesIO
            from pptx import Presentation

            # .ppt 格式需要先转换，这里只支持 .pptx
            if file_ext == '.ppt':
                return "暂不支持 .ppt 格式的内容预览，请转换为 .pptx 格式"

            # 从字节流打开演示文稿
            prs = Presentation(BytesIO(file_bytes))
            text_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                slide_texts.append(f"--- 幻灯片 {slide_num} ---")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # 处理表格
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_texts.append(" | ".join(row_text))

                if len(slide_texts) > 1:  # 除了标题外还有内容
                    text_content.append("\n".join(slide_texts))

            return "\n\n".join(text_content) if text_content else "PPT 文件内容为空"

        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return f"PPT 文件解析失败: {str(e)}"

    @classmethod
    def _extract_text_content(cls, file_bytes: bytes) -> str:
        """从文本文件提取内容"""
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']

            for encoding in encodings:
                try:
                    content = file_bytes.decode(encoding)
                    return content if content.strip() else "文本文件内容为空"
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用 utf-8 并忽略错误
            return file_bytes.decode('utf-8', errors='ignore')

        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return f"文本文件解析失败: {str(e)}"
